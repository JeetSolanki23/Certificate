from pptx import Presentation
import re

def extract_text_between_curly_braces(text):
    pattern = r'\{\{(.+?)\}\}'
    matches = re.findall(pattern, text)
    return matches

def main():
    presentation = Presentation('CERTIFICATE.pptx')
    user_inputs = {}

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):  # Check if shape has a text frame
                text = shape.text
                matches = extract_text_between_curly_braces(text)
                for match in matches:
                    if match not in user_inputs:
                        new_text = input(f"Enter value for '{match}': ")
                        user_inputs[match] = new_text

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):  # Check if shape has a text frame
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for match, new_text in user_inputs.items():
                            run.text = run.text.replace(f"{{{{{match}}}}}", new_text)

    updated_pptx_path = 'updated_presentation.pptx'
    presentation.save(updated_pptx_path)
    print(f"Presentation updated and saved as '{updated_pptx_path}'")

if __name__ == "__main__":
    main()
